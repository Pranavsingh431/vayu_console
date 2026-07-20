#!/usr/bin/env python3
"""Build the ET AI Hackathon 2.0 submission deck.

One slide definition, three artefacts: `final-presentation.pptx`,
`final-presentation.pdf`, and `presentation-content.md`. Holding the content in
one place is the point — a deck whose PDF and PPTX disagree about a number is
worse than having only one of them.

The PDF is rendered by printing HTML through headless Chromium (Playwright is
already a dev dependency for the E2E suite) rather than converting the PPTX,
because no converter is installed and the HTML path matches the product's own
palette exactly.

    python3 submission/build_presentation.py

Every figure in here is verified against the repository or the live API. See
`presentation-content.md` for the provenance table.
"""

from __future__ import annotations

import html
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SUBMISSION = ROOT / "submission"
SHOTS = ROOT / "docs" / "assets" / "screenshots"

# The product palette, verbatim from apps/web/src/app/globals.css.
BLACK = "000000"
SURFACE = "0A0A0A"
RAISED = "111111"
BORDER = "1C1C1C"
WHITE = "FFFFFF"
DIM = "A1A1AA"
FAINT = "71717A"
GHOST = "52525B"
GREEN = "22C55E"
AMBER = "EAB308"
RED = "EF4444"

W, H = Inches(13.333), Inches(7.5)  # 16:9

#: Separator between items on one architecture layer. Runs of spaces collapse in
#: HTML, so the two renderers need a character, not whitespace.
LAYER_SEP = "   ·   "


def rgb(hex_: str) -> RGBColor:
    return RGBColor.from_string(hex_)


# ---------------------------------------------------------------------------
# Slide content
# ---------------------------------------------------------------------------
#
# `kind` drives layout. `notes` become PowerPoint speaker notes and the "Speaker
# notes" block in the markdown.

SLIDES: list[dict] = [
    {
        "kind": "title",
        "eyebrow": "ET AI HACKATHON 2.0",
        "title": "Vayu Console",
        "subtitle": "From pollution signals to defensible action",
        "meta": [
            "Urban Air Quality Decision Intelligence",
            "Team: singhpranav431",
        ],
        "image": "02-console-diwali.png",
        "notes": (
            "Open flat. Do not read the slide. Say: this is a decision support console for "
            "municipal air quality officers, and the one thing that makes it different is that "
            "it refuses to answer when the evidence cannot support an answer. Then move on "
            "quickly, the problem slide is where the argument starts."
        ),
    },
    {
        "kind": "problem",
        "eyebrow": "THE PROBLEM",
        "title": "AQI tells cities how bad the air is.\nIt does not tell them what action the evidence justifies.",
        "left_label": "WHAT THE DASHBOARD SAYS",
        "left_value": "Severe",
        "left_sub": "1288 µg/m³ PM2.5",
        "right_label": "WHAT THE OFFICER STILL HAS TO ANSWER",
        "questions": [
            "Why is this happening?",
            "What evidence do we actually have?",
            "What action is justified right now?",
            "How confident should I be?",
            "Can I defend this decision tomorrow?",
        ],
        "footer": (
            "Cities have the monitoring data. What is missing is the reasoning layer between a "
            "reading and a defensible decision, and officers assemble it by hand."
        ),
        "notes": (
            "The audience already believes the first half. Spend your time on the five questions "
            "on the right. Land this line: none of these are measurement problems, they are "
            "reasoning problems, and that is the gap the product sits in."
        ),
    },
    {
        "kind": "flow4",
        "eyebrow": "THE SOLUTION",
        "title": "Vayu Console turns environmental evidence into transparent operational decisions.",
        "steps": [
            ("Situation", "What is happening?", "Observed pollution conditions for a station-hour."),
            ("Evidence", "What does it support?", "Three competing hypotheses, rated independently."),
            ("Decision", "What is justified?", "Deterministic rules produce recommendations."),
            ("Explanation", "Why this?", "Every recommendation traces back to its evidence."),
        ],
        "footer": "The layout of the console is this sequence, left to right. The interface is the workflow.",
        "notes": (
            "Four words. Say that evidence comes before decision structurally, not just visually: "
            "the Decision Engine can only read the Evidence Engine's output, so no rule can reach "
            "a recommendation the evidence does not carry."
        ),
    },
    {
        "kind": "architecture",
        "eyebrow": "HOW IT WORKS",
        "title": "Architecture",
        "layers": [
            ("DATA SOURCES", ["OpenAQ S3 archive (CPCB station readings)", "NASA FIRMS (VIIRS active fire)", "Open-Meteo (wind, boundary layer)"]),
            ("INGESTION & VALIDATION", ["Sentinel rejection · UTC normalisation · station deduplication"]),
            ("STORAGE", ["Supabase PostgreSQL + PostGIS"]),
            ("EVIDENCE ENGINE", ["Biomass / Fire", "Traffic", "Industrial"]),
            ("DECISION ENGINE", ["Deterministic rules → policies → recommendations"]),
            ("OPERATIONS CONSOLE", ["FastAPI  →  Next.js  →  Officer"]),
        ],
        "footer": "Three ingested sources. No LLM anywhere in the decision path.",
        "notes": (
            "Twenty seconds unless asked. The only detail worth volunteering: readings come from "
            "OpenAQ's S3 archive rather than its live API, because the two disagree about which "
            "Delhi stations have history. That is on the technical differentiation slide too."
        ),
    },
    {
        "kind": "evidence",
        "eyebrow": "THE EVIDENCE ENGINE",
        "title": "We evaluate hypotheses. We do not fabricate attribution.",
        "cards": [
            ("Fire / Biomass", "★★★☆☆", "MODERATE", "High-quality data", GREEN),
            ("Traffic", "★☆☆☆☆", "VERY WEAK", "Fair data", AMBER),
            ("Industry / Power", "★☆☆☆☆", "VERY WEAK", "Poor data", AMBER),
        ],
        "caption": "Actual output for the Diwali 2019 worked example.",
        "exposes": [
            "Supporting observations",
            "Contradicting observations",
            "Assumptions",
            "Evidence quality",
            "Limitations",
            "Historical validation",
        ],
        "callout": (
            "Evidence strengths are independent. They are not probabilities, they do not sum to "
            "100%, and there is no pie chart anywhere in the product."
        ),
        "notes": (
            "The line that lands: these do not add up to 100% and that is not an oversight. The "
            "hypotheses are not mutually exclusive. On a Diwali night in stubble season, fire and "
            "traffic evidence are both genuinely present. Any chart implying shares would be a lie "
            "about what we can know."
        ),
    },
    {
        "kind": "honesty",
        "eyebrow": "SCIENTIFIC HONESTY",
        "title": "Our first source classifier hit 98.6% accuracy.\nWe deleted it.",
        "chain": ["Features", "Rule-generated labels", "Random forest", "98.6% accuracy"],
        "verdict": "The model had learned the labelling rule, not the atmosphere.",
        "detail": (
            "No labelled dataset of pollution source exists for Delhi, because producing one would "
            "itself require apportionment. The only labels available were ones we invented from a "
            "rule over our own features. Training on that reproduces the rule."
        ),
        "measured": "13,685 real station-hours · 5-fold CV · accuracy 0.9859 ± 0.0035 · macro F1 0.9851",
        "rejected": "Fabricated source accuracy",
        "built": "Evidence-based hypothesis screening",
        "closer": "Statistical machinery cannot manufacture ground truth that does not exist.",
        "notes": (
            "Do not rush this slide, it is the most memorable thing you will say. The key beat: "
            "the danger was not that it failed, it is that it succeeded spectacularly. Every "
            "honesty instrument we could have pointed at it, SHAP, conformal prediction, "
            "calibration curves, would have corroborated the fabrication and made it harder to "
            "detect, not easier."
        ),
    },
    {
        "kind": "validation",
        "eyebrow": "HISTORICAL VALIDATION",
        "title": "Each module declares the intervention that would reject it.",
        "rows": [
            {
                "event": "COVID lockdown 2020",
                "tests": "Traffic module",
                "status": "ACCEPTED",
                "tone": GREEN,
                "detail": "47 stations · 901,160 rows.  NO₂ −54.4%, SO₂ −3.7%.  NO₂/SO₂ ratio 2.82 → 1.34.  Likelihood ratio 2.11 (weak).",
            },
            {
                "event": "Diwali 2019",
                "tests": "Fire discriminant validity",
                "status": "PENDING",
                "tone": FAINT,
                "detail": "1,604 VIIRS detections that day, so the confound is real. VIIRS overpasses Delhi 12:00–14:00 and 01:00–03:00 IST, outside the firework window. Structural argument, not yet a result.",
            },
            {
                "event": "Odd-Even II, Apr 2016",
                "tests": "Vehicular",
                "status": "DATA INGESTED — VALIDATION PENDING",
                "tone": FAINT,
                "detail": "The only unconfounded vehicle window in the record: no stubble, no winter inversion. Weak treatment on few stations.",
            },
        ],
        "callout": (
            "One test has been run to completion. The engine reports the rest as pending at "
            "/evidence/history, and the console shows them as pending."
        ),
        "notes": (
            "Be first to say that only one test is finished. It is a stronger position than "
            "pretending three are. The COVID beat: traffic went to zero by government decree, "
            "which is an experiment nobody could run deliberately, and it could have failed. Had "
            "NO₂ not fallen relative to SO₂ we would have deleted the traffic module."
        ),
    },
    {
        "kind": "chain",
        "eyebrow": "FROM EVIDENCE TO ACTION",
        "title": "Every recommendation has to earn its way in.",
        "generic": ["EVIDENCE", "RULE", "POLICY", "RECOMMENDATION"],
        "concrete": [
            "Biomass evidence\nMODERATE",
            "FIRE_002",
            "PUBLIC_HEALTH",
            "Issue a public health\nadvisory for the\naffected zone",
        ],
        "points": [
            "No LLM decides enforcement actions.",
            "The decision engine is deterministic: the same evidence returns the same recommendation, every time.",
            "A rule fires on the presence of evidence, never on its absence.",
        ],
        "notes": (
            "If asked why not an LLM: a recommendation an officer defends in a regulatory meeting "
            "has to be reproducible, and auditability is impossible without it. Also, fluency is "
            "exactly the failure mode we are avoiding. The hard part is not producing plausible "
            "explanations, it is refusing to produce them when the evidence is absent."
        ),
    },
    {
        "kind": "shot",
        "eyebrow": "CHALLENGE THE RECOMMENDATION",
        "title": "Do not trust the recommendation. Challenge it.",
        "image": "03-challenge-modal.png",
        "side": [
            "Evidence used",
            "Rule triggered",
            "Policy applied",
            "Supporting observations",
            "Contradicting observations",
            "Assumptions",
            "Historical validation",
            "Known limitations",
        ],
        "callout": "Every recommendation is auditable back to its evidence.",
        "notes": (
            "Do this live if the demo is working. The framing that lands: this is not built to "
            "convince the officer the system is right, it is built so someone can prove it wrong. "
            "An officer needs the counter-argument before someone else in the meeting supplies it, "
            "which is why contradicting evidence is never hidden."
        ),
    },
    {
        "kind": "product",
        "eyebrow": "THE PRODUCT",
        "title": "Historical Replay — Diwali 2019",
        "image": "02-console-diwali.png",
        "steps": [
            "Situation",
            "Transport geometry",
            "Evidence",
            "Recommendations",
            "Challenge",
            "Historical validation",
        ],
        "callout": (
            "Every scenario is a reconstruction from archived observations. The console says so on "
            "every screen. Nothing here is a live feed."
        ),
        "notes": (
            "Volunteer the replay point before anyone wonders about it. It reads as rigour when you "
            "say it first and as evasion when you are asked. One non-scrolling screen, designed for "
            "a control room at 2am."
        ),
    },
    {
        "kind": "differentiation",
        "eyebrow": "TECHNICAL DIFFERENTIATION",
        "title": "Built for evidence integrity, not just model output.",
        "items": [
            ("Resilient data foundation", "OpenAQ's live API advertises Delhi stations it will not return history for. Coverage was probed against the raw S3 archive instead, the only source whose claimed coverage matches what it returns."),
            ("Data quality at the boundary", "Missing-data sentinels (−999, −9999, any negative concentration) are dropped at parse time, not filtered in queries, so one can never reach the database and be forgotten by a later aggregate."),
            ("Station identity", "One physical station appears under several OpenAQ location ids. Left unresolved this overstates the network and double-counts readings into city means."),
            ("Observation gaps", "A satellite non-observation is not a true non-detection. VIIRS passes over Delhi twice a day, so an absence of fire detections may be an absence of overpass, and the system reports it that way."),
            ("Auditable decisions", "Evidence → Rule → Policy → Recommendation, exposed in the interface rather than buried in a log."),
        ],
        "notes": (
            "Pick two, not five, unless the room is technical. The OpenAQ S3 discovery and the "
            "observation-gap distinction are the ones that signal you actually handled the data "
            "rather than downloading a clean CSV."
        ),
    },
    {
        "kind": "future",
        "eyebrow": "IMPACT & FUTURE",
        "title": "From a Delhi proof of concept to city-scale environmental decision support.",
        "now_label": "BUILT AND DEPLOYED — DELHI",
        "now": [
            "Evidence Engine, three hypotheses",
            "Deterministic Decision Engine",
            "Operations Console with full audit trace",
            "COVID lockdown validation, completed",
        ],
        "next_label": "NOT YET BUILT",
        "next": [
            "Road-density context for the traffic hypothesis",
            "Industrial proximity for the industrial hypothesis",
            "Expanded historical validation (Diwali, Odd-Even)",
            "Live ingestion scheduling",
            "Additional cities, each needing its own calibration",
            "Policy-specific integrations",
        ],
        "callout": (
            "The architecture ports to another city. The calibration does not — every likelihood "
            "ratio is measured against a local natural experiment, so scaling is a data problem, "
            "not a config change."
        ),
        "notes": (
            "Be explicit that the right column is not built. Claiming one-click multi-city would "
            "undercut the entire honesty argument you just made two slides ago. Say the "
            "calibration line out loud, it shows you understand your own limits."
        ),
    },
    {
        "kind": "closing",
        "quote": (
            "Most systems are designed to always give an answer.\n"
            "Vayu Console is designed to give an answer only when the evidence can support one, "
            "and to show exactly why."
        ),
        "title": "Vayu Console",
        "meta": [
            "Team: singhpranav431",
            "ET AI Hackathon 2.0",
        ],
        "links": [
            "Console   vayu-console-web.vercel.app/console",
            "API       vayu-console-api.onrender.com",
            "Code      github.com/Pranavsingh431/vayu_console",
        ],
        "notes": (
            "Land the quote and stop talking. Do not add a summary after it. If there are "
            "questions, the judge FAQ in submission/judge-faq.md covers the hard ones, including "
            "why not an LLM, why not a classifier, and whether this is source apportionment."
        ),
    },
]


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------


def bg(slide, colour: str = BLACK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(colour)


def box(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    size: int,
    colour: str = WHITE,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Helvetica Neue",
    spacing: float | None = None,
    line: float = 1.25,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, chunk in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = rgb(colour)
        if spacing is not None:
            # python-pptx has no character-spacing API; set it on the run's XML.
            r.font._rPr.set("spc", str(int(spacing * 100)))
    return tb


def panel(slide, x, y, w, h, fill_c: str = SURFACE, line_c: str = BORDER, line_w: float = 1.0):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.03  # barely rounded; operations software, not a toy
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_c)
    shp.line.color.rgb = rgb(line_c)
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def rule(slide, x, y, w, colour: str = BORDER, thickness: float = 1.0):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(thickness))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(colour)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def eyebrow(slide, text: str, y=Inches(0.55)):
    box(slide, Inches(0.85), y, Inches(11), Inches(0.3), text, 11, AMBER, True, spacing=2.4)


def heading(slide, text: str, y=Inches(1.0), size=30, width=11.6):
    return box(slide, Inches(0.85), y, Inches(width), Inches(1.5), text, size, WHITE, True, line=1.18)


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def picture_fit(slide, path: Path, x, y, w, h):
    """Insert an image scaled to fit inside a box, centred, preserving aspect."""
    from PIL import Image  # noqa: PLC0415

    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 2880, 1800
    scale = min(w / iw, h / ih)
    dw, dh = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(str(path), int(x + (w - dw) / 2), int(y + (h - dh) / 2), dw, dh)


# ---------------------------------------------------------------------------
# PPTX rendering
# ---------------------------------------------------------------------------


def build_pptx(out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        bg(slide)
        kind = s["kind"]

        if kind == "title":
            img = SHOTS / s["image"]
            if img.exists():
                pic = picture_fit(slide, img, Inches(6.6), Inches(1.7), Inches(6.4), Inches(4.0))
                # Sit it behind the text block, quietly.
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
            eyebrow(slide, s["eyebrow"], Inches(1.5))
            box(slide, Inches(0.85), Inches(2.0), Inches(6.0), Inches(1.2), s["title"], 54, WHITE, True)
            box(slide, Inches(0.85), Inches(3.3), Inches(5.6), Inches(1.0), s["subtitle"], 20, DIM, line=1.35)
            rule(slide, Inches(0.85), Inches(4.5), Inches(1.2), AMBER, 2)
            box(slide, Inches(0.85), Inches(4.9), Inches(5.6), Inches(1.0), "\n".join(s["meta"]), 13, FAINT, line=1.6)

        elif kind == "problem":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 26)
            panel(slide, Inches(0.85), Inches(2.75), Inches(3.6), Inches(2.5))
            box(slide, Inches(1.15), Inches(3.05), Inches(3.0), Inches(0.3), s["left_label"], 9.5, FAINT, True, spacing=1.6)
            box(slide, Inches(1.15), Inches(3.5), Inches(3.0), Inches(0.7), s["left_value"], 34, RED, True)
            box(slide, Inches(1.15), Inches(4.3), Inches(3.0), Inches(0.4), s["left_sub"], 14, WHITE, font="Menlo")
            box(slide, Inches(4.75), Inches(3.05), Inches(0.4), Inches(0.4), "→", 20, GHOST)
            box(slide, Inches(5.5), Inches(3.05), Inches(6.5), Inches(0.3), s["right_label"], 9.5, FAINT, True, spacing=1.6)
            for i, q in enumerate(s["questions"]):
                y = Inches(3.55 + i * 0.42)
                box(slide, Inches(5.5), y, Inches(0.25), Inches(0.3), "·", 14, AMBER, True)
                box(slide, Inches(5.8), y, Inches(6.4), Inches(0.35), q, 14, WHITE)
            rule(slide, Inches(0.85), Inches(5.75), Inches(11.6))
            box(slide, Inches(0.85), Inches(5.95), Inches(11.6), Inches(0.8), s["footer"], 12.5, FAINT, line=1.5)

        elif kind == "flow4":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 26)
            cw, gap = Inches(2.72), Inches(0.28)
            for i, (name, q, desc) in enumerate(s["steps"]):
                x = Inches(0.85) + i * (cw + gap)
                panel(slide, x, Inches(2.65), cw, Inches(2.45))
                box(slide, x + Inches(0.28), Inches(2.95), Inches(0.5), Inches(0.3), f"0{i+1}", 10, GHOST, True, font="Menlo")
                box(slide, x + Inches(0.28), Inches(3.3), cw - Inches(0.56), Inches(0.45), name, 21, WHITE, True)
                box(slide, x + Inches(0.28), Inches(3.85), cw - Inches(0.56), Inches(0.35), q, 12, AMBER)
                box(slide, x + Inches(0.28), Inches(4.3), cw - Inches(0.56), Inches(0.8), desc, 11.5, DIM, line=1.4)
                if i < 3:
                    box(slide, x + cw + Inches(0.02), Inches(3.6), Inches(0.3), Inches(0.3), "→", 15, GHOST)
            rule(slide, Inches(0.85), Inches(5.6), Inches(11.6))
            box(slide, Inches(0.85), Inches(5.8), Inches(11.6), Inches(0.6), s["footer"], 12.5, FAINT)

        elif kind == "architecture":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 28)
            y = Inches(1.95)
            for i, (label, rows) in enumerate(s["layers"]):
                hgt = Inches(0.62)
                accent = AMBER if label in ("EVIDENCE ENGINE", "DECISION ENGINE") else BORDER
                panel(slide, Inches(0.85), y, Inches(11.6), hgt, RAISED if accent == AMBER else SURFACE, accent)
                box(slide, Inches(1.15), y + Inches(0.11), Inches(3.0), Inches(0.4), label, 10.5,
                    AMBER if accent == AMBER else FAINT, True, spacing=1.6)
                box(slide, Inches(4.3), y + Inches(0.1), Inches(8.0), Inches(0.45),
                    LAYER_SEP.join(rows), 12.5, WHITE, line=1.2)
                y = y + hgt + Inches(0.14)
                if i < len(s["layers"]) - 1:
                    box(slide, Inches(6.55), y - Inches(0.16), Inches(0.3), Inches(0.2), "↓", 11, GHOST)
            box(slide, Inches(0.85), Inches(6.75), Inches(11.6), Inches(0.5), s["footer"], 12.5, FAINT)

        elif kind == "evidence":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 27)
            cw, gap = Inches(3.72), Inches(0.22)
            for i, (name, stars, strength, quality, tone) in enumerate(s["cards"]):
                x = Inches(0.85) + i * (cw + gap)
                panel(slide, x, Inches(2.5), cw, Inches(1.75))
                box(slide, x + Inches(0.3), Inches(2.75), cw - Inches(0.6), Inches(0.4), name, 17, WHITE, True)
                box(slide, x + Inches(0.3), Inches(3.25), cw - Inches(0.6), Inches(0.4), stars, 20, tone, font="Menlo")
                box(slide, x + Inches(0.3), Inches(3.75), cw - Inches(0.6), Inches(0.3), f"{strength}  ·  {quality.upper()}", 9.5, FAINT, True, spacing=1.2)
            box(slide, Inches(0.85), Inches(4.35), Inches(11.6), Inches(0.3), s["caption"], 10.5, GHOST)
            box(slide, Inches(0.85), Inches(4.85), Inches(3.4), Inches(0.3), "EACH HYPOTHESIS EXPOSES", 9.5, FAINT, True, spacing=1.6)
            for i, item in enumerate(s["exposes"]):
                col, row = i % 3, i // 3
                box(slide, Inches(0.85) + col * Inches(3.94), Inches(5.25) + row * Inches(0.36),
                    Inches(3.8), Inches(0.32), f"·  {item}", 12.5, DIM)
            panel(slide, Inches(0.85), Inches(6.15), Inches(11.6), Inches(0.72), RAISED, AMBER)
            box(slide, Inches(1.15), Inches(6.34), Inches(11.0), Inches(0.5), s["callout"], 12.5, WHITE, line=1.3)

        elif kind == "honesty":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 29)
            cw, gap = Inches(2.5), Inches(0.42)
            for i, step in enumerate(s["chain"]):
                x = Inches(0.85) + i * (cw + gap)
                last = i == len(s["chain"]) - 1
                panel(slide, x, Inches(2.55), cw, Inches(0.72), RAISED if last else SURFACE, RED if last else BORDER)
                box(slide, x, Inches(2.75), cw, Inches(0.4), step, 13, RED if last else DIM, last, align=PP_ALIGN.CENTER)
                if not last:
                    box(slide, x + cw + Inches(0.08), Inches(2.75), Inches(0.3), Inches(0.3), "→", 14, GHOST)
            box(slide, Inches(0.85), Inches(3.5), Inches(11.6), Inches(0.4), s["verdict"], 17, WHITE, True)
            box(slide, Inches(0.85), Inches(4.0), Inches(11.6), Inches(0.9), s["detail"], 12.5, DIM, line=1.45)
            box(slide, Inches(0.85), Inches(4.95), Inches(11.6), Inches(0.3), s["measured"], 11, GHOST, font="Menlo")
            panel(slide, Inches(0.85), Inches(5.45), Inches(5.6), Inches(0.85), SURFACE, RED)
            box(slide, Inches(1.15), Inches(5.6), Inches(5.0), Inches(0.28), "WE REJECTED", 9.5, RED, True, spacing=1.6)
            box(slide, Inches(1.15), Inches(5.9), Inches(5.0), Inches(0.3), s["rejected"], 14, WHITE)
            panel(slide, Inches(6.85), Inches(5.45), Inches(5.6), Inches(0.85), SURFACE, GREEN)
            box(slide, Inches(7.15), Inches(5.6), Inches(5.0), Inches(0.28), "WE BUILT", 9.5, GREEN, True, spacing=1.6)
            box(slide, Inches(7.15), Inches(5.9), Inches(5.0), Inches(0.3), s["built"], 14, WHITE)
            box(slide, Inches(0.85), Inches(6.6), Inches(11.6), Inches(0.4), s["closer"], 14, AMBER, True)

        elif kind == "validation":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 28)
            y = Inches(2.15)
            for r in s["rows"]:
                hgt = Inches(1.28)
                panel(slide, Inches(0.85), y, Inches(11.6), hgt, SURFACE, BORDER)
                rule(slide, Inches(0.85), y, Inches(0.05), r["tone"], 0)
                from pptx.enum.shapes import MSO_SHAPE

                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), y, Emu(30000), hgt)
                bar.fill.solid()
                bar.fill.fore_color.rgb = rgb(r["tone"])
                bar.line.fill.background()
                bar.shadow.inherit = False
                box(slide, Inches(1.2), y + Inches(0.2), Inches(3.4), Inches(0.4), r["event"], 16, WHITE, True)
                box(slide, Inches(1.2), y + Inches(0.62), Inches(3.4), Inches(0.3), r["tests"].upper(), 9.5, GHOST, True, spacing=1.4)
                box(slide, Inches(4.8), y + Inches(0.22), Inches(3.4), Inches(0.3), r["status"], 10.5, r["tone"], True, spacing=1.4)
                box(slide, Inches(4.8), y + Inches(0.6), Inches(7.3), Inches(0.6), r["detail"], 11.5, DIM, line=1.35)
                y = y + hgt + Inches(0.16)
            panel(slide, Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.72), RAISED, AMBER)
            box(slide, Inches(1.15), Inches(6.54), Inches(11.0), Inches(0.5), s["callout"], 12.5, WHITE, line=1.3)

        elif kind == "chain":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 28)
            cw, gap = Inches(2.62), Inches(0.36)
            for i, (g, c) in enumerate(zip(s["generic"], s["concrete"])):
                x = Inches(0.85) + i * (cw + gap)
                box(slide, x, Inches(2.35), cw, Inches(0.3), g, 9.5, FAINT, True, spacing=1.8, align=PP_ALIGN.CENTER)
                panel(slide, x, Inches(2.75), cw, Inches(1.5), RAISED if i == 3 else SURFACE, AMBER if i == 3 else BORDER)
                tb = box(slide, x + Inches(0.2), Inches(3.05), cw - Inches(0.4), Inches(1.0), c,
                         13 if i == 3 else 14, WHITE, True, align=PP_ALIGN.CENTER, line=1.3,
                         font="Menlo" if i in (1, 2) else "Helvetica Neue")
                tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if i < 3:
                    box(slide, x + cw + Inches(0.03), Inches(3.3), Inches(0.32), Inches(0.3), "→", 15, GHOST)
            rule(slide, Inches(0.85), Inches(4.75), Inches(11.6))
            for i, p in enumerate(s["points"]):
                box(slide, Inches(0.85), Inches(5.0) + i * Inches(0.52), Inches(0.3), Inches(0.35), "·", 15, AMBER, True)
                box(slide, Inches(1.2), Inches(5.0) + i * Inches(0.52), Inches(11.2), Inches(0.45), p, 14, WHITE, line=1.35)

        elif kind == "shot":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 29)
            img = SHOTS / s["image"]
            if img.exists():
                picture_fit(slide, img, Inches(0.85), Inches(2.2), Inches(8.2), Inches(4.5))
            box(slide, Inches(9.4), Inches(2.2), Inches(3.2), Inches(0.3), "THE MODAL EXPOSES", 9.5, FAINT, True, spacing=1.6)
            for i, item in enumerate(s["side"]):
                box(slide, Inches(9.4), Inches(2.62) + i * Inches(0.38), Inches(3.3), Inches(0.34), f"·  {item}", 12.5, DIM)
            panel(slide, Inches(9.4), Inches(5.85), Inches(3.1), Inches(0.85), RAISED, AMBER)
            box(slide, Inches(9.65), Inches(6.02), Inches(2.7), Inches(0.6), s["callout"], 12, WHITE, True, line=1.25)

        elif kind == "product":
            eyebrow(slide, s["eyebrow"])
            box(slide, Inches(0.85), Inches(0.95), Inches(8.0), Inches(0.5), s["title"], 24, WHITE, True)
            img = SHOTS / s["image"]
            if img.exists():
                picture_fit(slide, img, Inches(0.85), Inches(1.75), Inches(9.1), Inches(5.0))
            box(slide, Inches(10.3), Inches(1.85), Inches(2.4), Inches(0.3), "READING ORDER", 9.5, FAINT, True, spacing=1.6)
            for i, step in enumerate(s["steps"]):
                y = Inches(2.3) + i * Inches(0.46)
                box(slide, Inches(10.3), y, Inches(0.35), Inches(0.3), f"{i+1}", 12, AMBER, True, font="Menlo")
                box(slide, Inches(10.7), y, Inches(2.3), Inches(0.35), step, 13, WHITE)
            panel(slide, Inches(10.3), Inches(5.2), Inches(2.35), Inches(1.5), RAISED, AMBER)
            box(slide, Inches(10.5), Inches(5.4), Inches(1.95), Inches(1.2), s["callout"], 10, DIM, line=1.35)

        elif kind == "differentiation":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 28)
            y = Inches(2.1)
            for i, (head, body) in enumerate(s["items"]):
                box(slide, Inches(0.85), y, Inches(0.45), Inches(0.35), f"0{i+1}", 11, AMBER, True, font="Menlo")
                box(slide, Inches(1.45), y - Inches(0.03), Inches(3.5), Inches(0.4), head, 15, WHITE, True)
                box(slide, Inches(5.1), y - Inches(0.02), Inches(7.35), Inches(0.8), body, 11.5, DIM, line=1.4)
                y = y + Inches(0.95)
                if i < len(s["items"]) - 1:
                    rule(slide, Inches(0.85), y - Inches(0.2), Inches(11.6))

        elif kind == "future":
            eyebrow(slide, s["eyebrow"])
            heading(slide, s["title"], Inches(1.0), 27)
            panel(slide, Inches(0.85), Inches(2.4), Inches(5.6), Inches(2.9), SURFACE, GREEN)
            box(slide, Inches(1.15), Inches(2.65), Inches(5.0), Inches(0.3), s["now_label"], 9.5, GREEN, True, spacing=1.5)
            for i, item in enumerate(s["now"]):
                box(slide, Inches(1.15), Inches(3.1) + i * Inches(0.46), Inches(5.1), Inches(0.4), f"·  {item}", 12.5, WHITE, line=1.3)
            panel(slide, Inches(6.85), Inches(2.4), Inches(5.6), Inches(2.9), SURFACE, BORDER)
            box(slide, Inches(7.15), Inches(2.65), Inches(5.0), Inches(0.3), s["next_label"], 9.5, FAINT, True, spacing=1.5)
            for i, item in enumerate(s["next"]):
                box(slide, Inches(7.15), Inches(3.1) + i * Inches(0.37), Inches(5.1), Inches(0.35), f"·  {item}", 11.5, DIM, line=1.25)
            panel(slide, Inches(0.85), Inches(5.55), Inches(11.6), Inches(1.0), RAISED, AMBER)
            box(slide, Inches(1.15), Inches(5.78), Inches(11.0), Inches(0.7), s["callout"], 13, WHITE, line=1.35)

        elif kind == "closing":
            rule(slide, Inches(0.85), Inches(1.9), Inches(1.2), AMBER, 2)
            box(slide, Inches(0.85), Inches(2.3), Inches(11.0), Inches(2.0), s["quote"], 25, WHITE, True, line=1.35)
            rule(slide, Inches(0.85), Inches(4.7), Inches(11.6))
            box(slide, Inches(0.85), Inches(5.0), Inches(5.0), Inches(0.6), s["title"], 30, WHITE, True)
            box(slide, Inches(0.85), Inches(5.75), Inches(5.0), Inches(0.8), "\n".join(s["meta"]), 13, FAINT, line=1.5)
            box(slide, Inches(6.85), Inches(5.05), Inches(5.6), Inches(1.5), "\n".join(s["links"]), 11.5, DIM, font="Menlo", line=1.7)

        notes(slide, s.get("notes", ""))

    prs.save(out)


# ---------------------------------------------------------------------------
# HTML -> PDF
# ---------------------------------------------------------------------------


def esc(t: str) -> str:
    return html.escape(t).replace("\n", "<br>")


def build_html(img_base: str) -> str:
    """Render the same slide data as printable HTML."""
    css = f"""
    @page {{ size: 1280px 720px; margin: 0; }}
    * {{ margin:0; padding:0; box-sizing:border-box; }}
    body {{ background:#{BLACK}; font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;
            -webkit-print-color-adjust:exact; print-color-adjust:exact; }}
    /* Centre the body block vertically. Content-light slides otherwise stack at
       the top and leave a third of the frame empty, which reads as unfinished
       rather than as deliberate whitespace. */
    .s {{ width:1280px; height:720px; background:#{BLACK}; color:#{WHITE};
          padding:52px 82px; position:relative; page-break-after:always; overflow:hidden;
          display:flex; flex-direction:column; justify-content:center; }}
    .s:last-child {{ page-break-after:auto; }}
    .s.top {{ justify-content:flex-start; }}
    .eb {{ font-size:11px; letter-spacing:.24em; color:#{AMBER}; font-weight:700; margin-bottom:18px; }}
    h1 {{ font-size:30px; line-height:1.2; font-weight:700; letter-spacing:-.01em; }}
    h1.big {{ font-size:56px; }}
    .p {{ background:#{SURFACE}; border:1px solid #{BORDER}; border-radius:6px; }}
    .p.am {{ background:#{RAISED}; border-color:#{AMBER}; }}
    .p.gr {{ border-color:#{GREEN}; }} .p.rd {{ border-color:#{RED}; }}
    .lbl {{ font-size:9.5px; letter-spacing:.16em; color:#{FAINT}; font-weight:700; }}
    .dim {{ color:#{DIM}; }} .faint {{ color:#{FAINT}; }} .ghost {{ color:#{GHOST}; }}
    .mono {{ font-family:Menlo,monospace; }}
    .row {{ display:flex; gap:14px; }}
    .arrow {{ color:#{GHOST}; font-size:15px; align-self:center; }}
    hr {{ border:0; border-top:1px solid #{BORDER}; }}
    img {{ display:block; border-radius:4px; border:1px solid #{BORDER}; }}
    .cal {{ background:#{RAISED}; border:1px solid #{AMBER}; border-radius:6px; padding:14px 20px;
            font-size:12.5px; line-height:1.4; }}
    """
    out = [f"<!doctype html><meta charset='utf-8'><style>{css}</style>"]

    for s in SLIDES:
        k = s["kind"]
        # Image-led slides keep their natural top-down flow; everything else is
        # centred so the frame reads as balanced at any content length.
        out.append(f"<div class='s{' top' if k in ('shot', 'product') else ''}'>")

        if k == "title":
            out.append(
                f"<img src='{img_base}/{s['image']}' style='position:absolute;right:0;top:150px;"
                f"width:620px;opacity:.55'>"
                f"<div style='position:relative'>"
                f"<div class='eb'>{esc(s['eyebrow'])}</div>"
                f"<h1 class='big'>{esc(s['title'])}</h1>"
                f"<div style='font-size:20px;color:#{DIM};margin-top:16px;max-width:520px;line-height:1.35'>{esc(s['subtitle'])}</div>"
                f"<div style='width:64px;height:2px;background:#{AMBER};margin:34px 0 18px'></div>"
                f"<div style='font-size:13px;color:#{FAINT};line-height:1.75'>{esc(chr(10).join(s['meta']))}</div></div>"
            )

        elif k == "problem":
            qs = "".join(
                f"<div style='display:flex;gap:12px;margin-bottom:11px'><span style='color:#{AMBER}'>·</span>"
                f"<span style='font-size:14px'>{esc(q)}</span></div>"
                for q in s["questions"]
            )
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:26px'>{esc(s['title'])}</h1>"
                f"<div class='row' style='margin-top:34px;align-items:stretch'>"
                f"<div class='p' style='width:340px;padding:24px'>"
                f"<div class='lbl'>{esc(s['left_label'])}</div>"
                f"<div style='font-size:38px;font-weight:700;color:#{RED};margin-top:14px'>{esc(s['left_value'])}</div>"
                f"<div class='mono' style='font-size:14px;margin-top:8px'>{esc(s['left_sub'])}</div></div>"
                f"<div class='arrow'>→</div>"
                f"<div style='flex:1;padding:22px 8px'><div class='lbl' style='margin-bottom:16px'>{esc(s['right_label'])}</div>{qs}</div>"
                f"</div><hr style='margin-top:26px'>"
                f"<div class='faint' style='font-size:12.5px;margin-top:16px;line-height:1.5'>{esc(s['footer'])}</div>"
            )

        elif k == "flow4":
            cards = ""
            for i, (n, q, d) in enumerate(s["steps"]):
                cards += (
                    f"<div class='p' style='flex:1;padding:22px'>"
                    f"<div class='mono ghost' style='font-size:10px'>0{i+1}</div>"
                    f"<div style='font-size:21px;font-weight:700;margin-top:12px'>{esc(n)}</div>"
                    f"<div style='font-size:12px;color:#{AMBER};margin-top:8px'>{esc(q)}</div>"
                    f"<div class='dim' style='font-size:11.5px;margin-top:12px;line-height:1.4'>{esc(d)}</div></div>"
                )
                if i < 3:
                    cards += "<div class='arrow'>→</div>"
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:26px'>{esc(s['title'])}</h1>"
                f"<div class='row' style='margin-top:34px'>{cards}</div><hr style='margin-top:30px'>"
                f"<div class='faint' style='font-size:12.5px;margin-top:16px'>{esc(s['footer'])}</div>"
            )

        elif k == "architecture":
            rows = ""
            for i, (label, items) in enumerate(s["layers"]):
                am = label in ("EVIDENCE ENGINE", "DECISION ENGINE")
                rows += (
                    f"<div class='p {'am' if am else ''}' style='padding:12px 22px;display:flex;align-items:center;gap:28px'>"
                    f"<div class='lbl' style='width:210px;{f'color:#{AMBER}' if am else ''}'>{esc(label)}</div>"
                    f"<div style='font-size:12.5px'>{esc(LAYER_SEP.join(items))}</div></div>"
                )
                if i < len(s["layers"]) - 1:
                    rows += f"<div style='text-align:center;color:#{GHOST};font-size:11px;margin:5px 0'>↓</div>"
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:28px'>{esc(s['title'])}</h1>"
                f"<div style='margin-top:24px'>{rows}</div>"
                f"<div class='faint' style='font-size:12.5px;margin-top:20px'>{esc(s['footer'])}</div>"
            )

        elif k == "evidence":
            cards = "".join(
                f"<div class='p' style='flex:1;padding:22px'>"
                f"<div style='font-size:17px;font-weight:700'>{esc(n)}</div>"
                f"<div class='mono' style='font-size:20px;color:#{tone};margin-top:12px'>{esc(st)}</div>"
                f"<div class='lbl' style='margin-top:12px'>{esc(sg)} &nbsp;·&nbsp; {esc(q.upper())}</div></div>"
                for n, st, sg, q, tone in s["cards"]
            )
            ex = "".join(
                f"<div class='dim' style='width:33%;font-size:12.5px;margin-bottom:9px'>·&nbsp; {esc(i)}</div>"
                for i in s["exposes"]
            )
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:27px'>{esc(s['title'])}</h1>"
                f"<div class='row' style='margin-top:26px'>{cards}</div>"
                f"<div class='ghost' style='font-size:10.5px;margin-top:12px'>{esc(s['caption'])}</div>"
                f"<div class='lbl' style='margin-top:22px;margin-bottom:12px'>EACH HYPOTHESIS EXPOSES</div>"
                f"<div style='display:flex;flex-wrap:wrap'>{ex}</div>"
                f"<div class='cal' style='margin-top:14px'>{esc(s['callout'])}</div>"
            )

        elif k == "honesty":
            ch = ""
            for i, step in enumerate(s["chain"]):
                last = i == len(s["chain"]) - 1
                ch += (
                    f"<div class='p {'rd' if last else ''}' style='flex:1;padding:14px;text-align:center;"
                    f"font-size:13px;{f'color:#{RED};font-weight:700;background:#{RAISED}' if last else f'color:#{DIM}'}'>{esc(step)}</div>"
                )
                if not last:
                    ch += "<div class='arrow'>→</div>"
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:29px'>{esc(s['title'])}</h1>"
                f"<div class='row' style='margin-top:24px'>{ch}</div>"
                f"<div style='font-size:17px;font-weight:700;margin-top:22px'>{esc(s['verdict'])}</div>"
                f"<div class='dim' style='font-size:12.5px;margin-top:12px;line-height:1.45;max-width:1000px'>{esc(s['detail'])}</div>"
                f"<div class='mono ghost' style='font-size:11px;margin-top:14px'>{esc(s['measured'])}</div>"
                f"<div class='row' style='margin-top:18px'>"
                f"<div class='p rd' style='flex:1;padding:14px 20px'><div class='lbl' style='color:#{RED}'>WE REJECTED</div>"
                f"<div style='font-size:14px;margin-top:7px'>{esc(s['rejected'])}</div></div>"
                f"<div class='p gr' style='flex:1;padding:14px 20px'><div class='lbl' style='color:#{GREEN}'>WE BUILT</div>"
                f"<div style='font-size:14px;margin-top:7px'>{esc(s['built'])}</div></div></div>"
                f"<div style='font-size:14px;font-weight:700;color:#{AMBER};margin-top:20px'>{esc(s['closer'])}</div>"
            )

        elif k == "validation":
            rows = "".join(
                f"<div class='p' style='padding:16px 22px;margin-bottom:12px;display:flex;gap:32px;"
                f"border-left:3px solid #{r['tone']}'>"
                f"<div style='width:250px'><div style='font-size:16px;font-weight:700'>{esc(r['event'])}</div>"
                f"<div class='lbl' style='margin-top:7px'>{esc(r['tests'].upper())}</div></div>"
                f"<div style='flex:1'><div class='lbl' style='color:#{r['tone']}'>{esc(r['status'])}</div>"
                f"<div class='dim' style='font-size:11.5px;margin-top:8px;line-height:1.4'>{esc(r['detail'])}</div></div></div>"
                for r in s["rows"]
            )
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:28px'>{esc(s['title'])}</h1>"
                f"<div style='margin-top:22px'>{rows}</div>"
                f"<div class='cal'>{esc(s['callout'])}</div>"
            )

        elif k == "chain":
            ch = ""
            for i, (g, c) in enumerate(zip(s["generic"], s["concrete"])):
                mono = "mono" if i in (1, 2) else ""
                ch += (
                    f"<div style='flex:1'><div class='lbl' style='text-align:center;margin-bottom:10px'>{esc(g)}</div>"
                    f"<div class='p {'am' if i == 3 else ''}' style='padding:20px 14px;height:112px;display:flex;"
                    f"align-items:center;justify-content:center;text-align:center'>"
                    f"<div class='{mono}' style='font-size:{13 if i == 3 else 14}px;font-weight:700;line-height:1.35'>{esc(c)}</div></div></div>"
                )
                if i < 3:
                    ch += "<div class='arrow' style='padding-top:26px'>→</div>"
            pts = "".join(
                f"<div style='display:flex;gap:12px;margin-bottom:11px'><span style='color:#{AMBER}'>·</span>"
                f"<span style='font-size:14px'>{esc(p)}</span></div>"
                for p in s["points"]
            )
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:28px'>{esc(s['title'])}</h1>"
                f"<div class='row' style='margin-top:26px'>{ch}</div><hr style='margin:28px 0 20px'>{pts}"
            )

        elif k == "shot":
            side = "".join(
                f"<div class='dim' style='font-size:12.5px;margin-bottom:10px'>·&nbsp; {esc(i)}</div>"
                for i in s["side"]
            )
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:29px'>{esc(s['title'])}</h1>"
                f"<div class='row' style='margin-top:24px'>"
                f"<img src='{img_base}/{s['image']}' style='width:790px'>"
                f"<div style='flex:1'><div class='lbl' style='margin-bottom:14px'>THE MODAL EXPOSES</div>{side}"
                f"<div class='cal' style='margin-top:18px;font-size:12px;font-weight:700'>{esc(s['callout'])}</div></div></div>"
            )

        elif k == "product":
            steps = "".join(
                f"<div style='display:flex;gap:12px;margin-bottom:12px'>"
                f"<span class='mono' style='color:#{AMBER};font-size:12px'>{i+1}</span>"
                f"<span style='font-size:13px'>{esc(st)}</span></div>"
                for i, st in enumerate(s["steps"])
            )
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div>"
                f"<div style='font-size:24px;font-weight:700'>{esc(s['title'])}</div>"
                f"<div class='row' style='margin-top:22px'>"
                f"<img src='{img_base}/{s['image']}' style='width:880px'>"
                f"<div style='flex:1'><div class='lbl' style='margin-bottom:14px'>READING ORDER</div>{steps}"
                f"<div class='cal' style='margin-top:20px;font-size:10.5px'>{esc(s['callout'])}</div></div></div>"
            )

        elif k == "differentiation":
            items = ""
            for i, (h_, b) in enumerate(s["items"]):
                items += (
                    f"<div style='display:flex;gap:22px;padding:14px 0;"
                    f"{'border-bottom:1px solid #' + BORDER if i < len(s['items']) - 1 else ''}'>"
                    f"<div class='mono' style='color:#{AMBER};font-size:11px;width:24px;padding-top:3px'>0{i+1}</div>"
                    f"<div style='width:330px;font-size:15px;font-weight:700'>{esc(h_)}</div>"
                    f"<div class='dim' style='flex:1;font-size:11.5px;line-height:1.45'>{esc(b)}</div></div>"
                )
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:28px'>{esc(s['title'])}</h1>"
                f"<div style='margin-top:14px'>{items}</div>"
            )

        elif k == "future":
            now = "".join(f"<div style='font-size:12.5px;margin-bottom:11px'>·&nbsp; {esc(i)}</div>" for i in s["now"])
            nxt = "".join(f"<div class='dim' style='font-size:11.5px;margin-bottom:9px'>·&nbsp; {esc(i)}</div>" for i in s["next"])
            out.append(
                f"<div class='eb'>{esc(s['eyebrow'])}</div><h1 style='font-size:27px'>{esc(s['title'])}</h1>"
                f"<div class='row' style='margin-top:26px;align-items:stretch'>"
                f"<div class='p gr' style='flex:1;padding:22px'><div class='lbl' style='color:#{GREEN};margin-bottom:16px'>{esc(s['now_label'])}</div>{now}</div>"
                f"<div class='p' style='flex:1;padding:22px'><div class='lbl' style='margin-bottom:16px'>{esc(s['next_label'])}</div>{nxt}</div>"
                f"</div><div class='cal' style='margin-top:20px;font-size:13px'>{esc(s['callout'])}</div>"
            )

        elif k == "closing":
            out.append(
                f"<div>"
                f"<div style='width:64px;height:2px;background:#{AMBER};margin-bottom:28px'></div>"
                f"<div style='font-size:25px;font-weight:700;line-height:1.4;max-width:1020px'>{esc(s['quote'])}</div>"
                f"<hr style='margin:44px 0 26px'>"
                f"<div class='row'><div style='flex:1'>"
                f"<div style='font-size:30px;font-weight:700'>{esc(s['title'])}</div>"
                f"<div class='faint' style='font-size:13px;margin-top:12px;line-height:1.6'>{esc(chr(10).join(s['meta']))}</div></div>"
                # white-space:pre keeps the column alignment in the link block;
                # HTML would otherwise collapse the padding spaces.
                f"<div class='mono dim' style='flex:1;font-size:11.5px;line-height:1.9;"
                f"padding-top:6px;white-space:pre'>{esc(chr(10).join(s['links']))}</div>"
                f"</div></div>"
            )

        out.append("</div>")

    return "".join(out)


def build_pdf(out: Path) -> bool:
    """Print the HTML deck to PDF with headless Chromium via Playwright."""
    tmp = Path(tempfile.mkdtemp())
    try:
        html_path = tmp / "deck.html"
        html_path.write_text(build_html(SHOTS.as_uri()), encoding="utf-8")

        from playwright.sync_api import sync_playwright  # noqa: PLC0415

        with sync_playwright() as pw:
            browser = pw.chromium.launch()
            page = browser.new_page()
            page.goto(html_path.as_uri())
            page.wait_for_timeout(1200)
            page.pdf(path=str(out), width="1280px", height="720px", print_background=True)
            browser.close()
        return True
    except Exception as exc:  # pragma: no cover
        print(f"  PDF step skipped: {exc}", file=sys.stderr)
        return False
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------


def build_markdown(out: Path) -> None:
    L: list[str] = [
        "# Vayu Console — presentation content",
        "",
        "Slide-by-slide content for `final-presentation.pptx` / `.pdf`, with speaker notes.",
        "",
        "Generated by `submission/build_presentation.py`, which holds the slide content as data",
        "and renders all three artefacts from it. Edit the script, not this file.",
        "",
        "**Team:** singhpranav431 · **Event:** ET AI Hackathon 2.0 · "
        f"**Slides:** {len(SLIDES)}",
        "",
        "---",
        "",
    ]

    for i, s in enumerate(SLIDES, 1):
        title = s.get("title", "Closing").replace("\n", " ")
        L += [f"## Slide {i} — {title}", ""]
        if s.get("eyebrow"):
            L += [f"**{s['eyebrow']}**", ""]

        k = s["kind"]
        if k == "title":
            L += [f"> {s['subtitle']}", "", *[f"- {m}" for m in s["meta"]], "",
                  f"Background: `docs/assets/screenshots/{s['image']}`", ""]
        elif k == "problem":
            L += [f"**{s['left_label']}:** {s['left_value']} ({s['left_sub']})", "",
                  f"**{s['right_label']}**", "", *[f"- {q}" for q in s["questions"]], "",
                  s["footer"], ""]
        elif k == "flow4":
            L += [f"- **{n}** — *{q}* {d}" for n, q, d in s["steps"]] + ["", s["footer"], ""]
        elif k == "architecture":
            L += ["```"]
            for j, (label, items) in enumerate(s["layers"]):
                L += [f"{label}", *[f"    {it}" for it in items]]
                if j < len(s["layers"]) - 1:
                    L += ["      ↓"]
            L += ["```", "", s["footer"], ""]
        elif k == "evidence":
            L += ["| Hypothesis | Strength | Rating | Data |", "| --- | --- | --- | --- |"]
            L += [f"| {n} | {st} | {sg} | {q} |" for n, st, sg, q, _ in s["cards"]]
            L += ["", f"*{s['caption']}*", "", "Each hypothesis exposes:", ""]
            L += [f"- {e}" for e in s["exposes"]] + ["", f"> {s['callout']}", ""]
        elif k == "honesty":
            L += [" → ".join(s["chain"]), "", f"**{s['verdict']}**", "", s["detail"], "",
                  f"`{s['measured']}`", "",
                  f"- **We rejected:** {s['rejected']}", f"- **We built:** {s['built']}", "",
                  f"> {s['closer']}", ""]
        elif k == "validation":
            L += ["| Experiment | Tests | Status | Detail |", "| --- | --- | --- | --- |"]
            L += [f"| {r['event']} | {r['tests']} | **{r['status']}** | {r['detail']} |" for r in s["rows"]]
            L += ["", f"> {s['callout']}", ""]
        elif k == "chain":
            L += [" → ".join(s["generic"]), "",
                  " → ".join(c.replace("\n", " ") for c in s["concrete"]), ""]
            L += [f"- {p}" for p in s["points"]] + [""]
        elif k == "shot":
            L += [f"Screenshot: `docs/assets/screenshots/{s['image']}`", "", "The modal exposes:", ""]
            L += [f"- {i_}" for i_ in s["side"]] + ["", f"> {s['callout']}", ""]
        elif k == "product":
            L += [f"Screenshot: `docs/assets/screenshots/{s['image']}`", "", "Reading order:", ""]
            L += [f"{j+1}. {st}" for j, st in enumerate(s["steps"])] + ["", f"> {s['callout']}", ""]
        elif k == "differentiation":
            L += [f"**{h_}** — {b}\n" for h_, b in s["items"]] + [""]
        elif k == "future":
            L += [f"**{s['now_label']}**", "", *[f"- {i_}" for i_ in s["now"]], "",
                  f"**{s['next_label']}**", "", *[f"- {i_}" for i_ in s["next"]], "",
                  f"> {s['callout']}", ""]
        elif k == "closing":
            L += [f"> {s['quote']}", "", *[f"- {m}" for m in s["meta"]], "",
                  "```", *s["links"], "```", ""]

        L += ["**Speaker notes.** " + s.get("notes", ""), "", "---", ""]

    L += [
        "## Verified figures and their sources",
        "",
        "Every quantitative claim in the deck, and where it comes from. Anything that could not",
        "be verified against the repository or the live API was left out.",
        "",
        "| Figure | Value | Source |",
        "| --- | --- | --- |",
        "| COVID stations | 47 | live `/evidence/history`, traffic module |",
        "| COVID rows | 901,160 | live `/evidence/history`, traffic module |",
        "| COVID NO₂ change | −54.4% | `apps/api/app/evidence/traffic/module.py` |",
        "| COVID SO₂ change | −3.7% | `apps/api/app/evidence/traffic/module.py` |",
        "| COVID NO₂/SO₂ ratio | 2.82 → 1.34 | `apps/api/app/evidence/traffic/module.py` |",
        "| Traffic likelihood ratio | 2.11 | `COVID_LIKELIHOOD_RATIO`, confirmed live |",
        "| Classifier station-hours | 13,685 | `docs/research/inference.md` §3 |",
        "| Classifier accuracy | 0.9859 ± 0.0035 | `docs/research/inference.md` §3 |",
        "| Classifier macro F1 | 0.9851 ± 0.0038 | `docs/research/inference.md` §3 |",
        "| Diwali VIIRS detections | 1,604 that day | `apps/api/app/evidence/biomass/module.py` |",
        "| VIIRS overpass windows | 12:00–14:00, 01:00–03:00 IST | biomass module / timeline |",
        "| Diwali example PM2.5 | 1288 µg/m³ | live `/evidence/example` |",
        "| Diwali example fire detections | 28 | live `/evidence/example` |",
        "| Sentinel values | −999, −9999, any negative | `app/data/sources/openaq_s3.py` |",
        "| Rule shown | FIRE_002 → PUBLIC_HEALTH | `app/decision/rules/rules.py` |",
        "",
        "## Claims deliberately excluded",
        "",
        "- **`data.gov.in` / CPCB as a direct ingested source.** Only three source modules exist",
        "  (`openaq_s3.py`, `firms.py`, `open_meteo.py`). CPCB readings arrive *via* the OpenAQ",
        "  archive, and the deck says so rather than listing a fourth pipeline that is not built.",
        "- **Diwali 2019 as a completed validation.** The engine reports it `pending`. An earlier",
        "  draft of the UI showed it as VERIFIED; that was corrected, and the deck says pending.",
        "- **Any per-event station count except COVID's 47**, which is the only one the live API",
        "  confirms.",
        "- **Any source percentage, probability, or causal attribution.** The system does not",
        "  produce them.",
        "- **Live operation.** The console replays reconstructed historical incidents.",
        "",
    ]
    out.write_text("\n".join(L), encoding="utf-8")


def main() -> int:
    SUBMISSION.mkdir(exist_ok=True)
    pptx_path = SUBMISSION / "final-presentation.pptx"
    pdf_path = SUBMISSION / "final-presentation.pdf"
    md_path = SUBMISSION / "presentation-content.md"

    missing = [s["image"] for s in SLIDES if s.get("image") and not (SHOTS / s["image"]).exists()]
    if missing:
        print(f"  missing screenshots: {missing}", file=sys.stderr)

    build_pptx(pptx_path)
    print(f"  {pptx_path.relative_to(ROOT)}  ({len(SLIDES)} slides)")

    if build_pdf(pdf_path):
        print(f"  {pdf_path.relative_to(ROOT)}")

    build_markdown(md_path)
    print(f"  {md_path.relative_to(ROOT)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
